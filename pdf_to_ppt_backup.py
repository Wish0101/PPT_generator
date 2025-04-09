import fitz  # PyMuPDF for normal PDFs
import pytesseract  # OCR for scanned PDFs
from pdf2image import convert_from_path
import os
import ollama
import time
from concurrent.futures import ThreadPoolExecutor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re

class PDFExtractor:
    def __init__(self, pdf_path):
        self.pdf_path = pdf_path
        self.text = None
        self.is_scanned = self.is_scanned_pdf(pdf_path)
        self.text = self.extract_text(pdf_path)

    def is_scanned_pdf(self, pdf_path):
        """Check if a PDF is scanned by analyzing the first page."""
        doc = fitz.open(pdf_path)
        first_page = doc[0]
        text = first_page.get_text("text").strip()
        return len(text) == 0  # If no text is found, it's likely a scanned PDF

    def extract_text_from_normal_pdf(self, pdf_path):
        """Extract text from a normal (non-scanned) PDF."""
        doc = fitz.open(pdf_path)
        text = "\n".join([page.get_text("text") for page in doc])
        return text

    def extract_text_from_scanned_pdf(self, pdf_path):
        """Extract text from a scanned PDF using OCR."""
        images = convert_from_path(pdf_path)
        text = "\n".join([pytesseract.image_to_string(img) for img in images])
        return text

    def extract_text(self, pdf_path):
        """Determine PDF type and extract text accordingly."""
        if self.is_scanned_pdf(pdf_path):
            print("Detected: Scanned PDF. Using OCR.")
            return self.extract_text_from_scanned_pdf(pdf_path)
        else:
            print("Detected: Normal PDF. Extracting text.")
            return self.extract_text_from_normal_pdf(pdf_path)


class PresentationGenerator:
    def __init__(self, text , model="gemmaModel" , theme=None):
        self.text = text
        self.model = model
        self.prs = Presentation()
        self.slide_width = self.prs.slide_width
        self.slide_height = self.prs.slide_height
        self.response_cache = {}

         # Default theme colors
        self.theme = {
            'header_bg': RGBColor(199, 198, 167),  # Red background
            'header_text': RGBColor(132, 174, 154),  # Yellow text
            'bullet_colors':[ RGBColor(199, 198, 167)],
            'slide_bg': RGBColor(132, 174, 154)  # Adding slide background color
        }

        if theme:
            self.theme.update(theme)
        
    def _apply_slide_background(self, slide):
        """
        Apply the theme background color to a slide
        
        Args:
            slide: The slide to apply the background color to
        """
        # Set background fill
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self.theme['slide_bg']
    def generate_response(self, prompt):
        """
        Generate a response using the AI model.
        
        Args:
            prompt (str): The prompt to send to the AI model
            
        Returns:
            str: The generated response text
        """
        # Check cache first
        if prompt in self.response_cache:
            return self.response_cache[prompt]
            
        try:
            # Send the prompt to the AI model and get the response
            response = ollama.chat(model=self.model, messages=[{"role": "user", "content": prompt}])
            
            # Extract the content from the message field
            if hasattr(response, 'message') and hasattr(response.message, 'content'):
                result = response.message.content
            elif isinstance(response, dict):
                if 'message' in response:
                    if isinstance(response['message'], dict) and 'content' in response['message']:
                        result = response['message']['content']
                    elif hasattr(response['message'], 'content'):
                        result = response['message'].content
                    else:
                        result = str(response['message'])
                else:
                    # Fallback: convert the entire response to string
                    result = str(response)
                    
                    # Try to extract content using regex
                    match = re.search(r"content='(.*?)'", result, re.DOTALL)
                    if match:
                        result = match.group(1)
            else:
                # Fallback for any other format
                result = str(response)
            
            # Clean up the response
            result = self._clean_response(result)
            
            # Cache the response
            self.response_cache[prompt] = result
            
            return result
        except Exception as e:
            print(f"Error generating response: {str(e)}")
            
            # Provide fallback content to keep the presentation generation going
            if "bullet points" in prompt.lower():
                return "• Important historical milestone\n• Key development in AI\n• Significant innovation"
            elif "one line" in prompt.lower() or "one clear" in prompt.lower():
                return "The history of artificial intelligence showcases humanity's quest to create machines that can think and learn."
            elif "conclusion" in prompt.lower():
                return "The history of AI demonstrates our continuous progress in creating intelligent systems. From early rule-based systems to modern neural networks, each advancement brings us closer to machines that can truly think."
            else:
                return "Content generation failed. Please try again with a different prompt."


    def _clean_response(self, text):
        """
        Clean up the response text by removing unwanted characters and formatting.
        
        Args:
            text (str): The text to clean
            
        Returns:
            str: The cleaned text
        """
        # Remove Markdown formatting if present
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # Remove bold formatting
        text = re.sub(r'\*(.*?)\*', r'\1', text)      # Remove italic formatting
        
        # Remove any lines that look like they're part of the API response structure
        text = re.sub(r'^.*?model=.*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'^.*?created_at=.*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'^.*?done=.*$', '', text, flags=re.MULTILINE)
        text = re.sub(r'^.*?message=.*$', '', text, flags=re.MULTILINE)
        
        # Remove any escaped newlines and clean up extra spaces
        text = text.replace('\\n', '\n')
        text = re.sub(r'\s{2,}', ' ', text)
        
        # Clean up lines
        lines = []
        for line in text.split('\n'):
            line = line.strip()
            if line and not line.startswith('Example:') and 'in real-world applications' not in line:
                lines.append(line)
        
        return '\n'.join(lines)

    def _extract_bullet_points(self, text):
        """
        Extract bullet points from text.
        
        Args:
            text (str): The text containing bullet points
            
        Returns:
            list: List of bullet points
        """
        # First, split by lines
        lines = text.split('\n')
        
        # Extract bullet points
        bullet_points = []
        for line in lines:
            line = line.strip()
            if line:
                # Remove bullet symbols and any other formatting
                line = re.sub(r'^[•\-\*]\s*', '', line)
                line = re.sub(r'^\d+\.\s*', '', line)
                
                # Remove any Markdown formatting
                line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)  # Remove bold formatting
                line = re.sub(r'\*(.*?)\*', r'\1', line)      # Remove italic formatting
                
                if line:
                    bullet_points.append(line)
        
        return bullet_points

    def _create_header(self, slide, title, top_margin=Inches(0)):
        """
        Create a standard header for slides.
        
        Args:
            slide: The slide to add the header to
            title (str): The title text for the header
            top_margin: The top margin for the header
            
        Returns:
            text_frame: The text frame of the header
        """
        rect_height = self.slide_height * 0.15
        rect_width = self.slide_width

        rectangle = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, top_margin, rect_width, rect_height
        )
        rectangle.fill.solid()
        rectangle.fill.fore_color.rgb = self.theme['header_bg']
        rectangle.line.fill.background()

        text_frame = rectangle.text_frame
        text_frame.text = title
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(24)
        run.font.color.rgb = self.theme['header_text']
        
        return text_frame

    def _create_text_box(self, slide, left=None, top=None, width=None, height=None, centered=False):
        """
        Create a standard text box for slides.
        
        Args:
            slide: The slide to add the text box to
            left: The left position of the text box
            top: The top position of the text box
            width: The width of the text box
            height: The height of the text box
            centered (bool): Whether to center the text box on the slide
            
        Returns:
            text_frame: The text frame of the text box
        """
        if centered:
            padding = Inches(0.1)
            top_padding = Inches(0.01)
            width = width or self.slide_width * 0.8
            height = height or Inches(2.5)
            left = (self.slide_width - width) / 2
            top = (self.slide_height - height) / 2
        else:
            left = left or Inches(1)
            top = top or Inches(2)
            width = width or (self.slide_width - Inches(2))
            height = height or (self.slide_height - Inches(3))
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        
        if centered:
            text_frame.margin_top = top_padding
            text_frame.margin_bottom = padding
            text_frame.margin_left = padding
            text_frame.margin_right = padding
        
        return text_frame

    def _add_paragraph(self, text_frame, text, color=None, size=Pt(14), level=0, is_bullet=False, bold=False):
        """
        Add a paragraph to a text frame.
        
        Args:
            text_frame: The text frame to add the paragraph to
            text (str): The text for the paragraph
            color: The color of the text
            size: The size of the text
            level (int): The indentation level
            is_bullet (bool): Whether to format as a bullet point
            bold (bool): Whether to make the text bold
            
        Returns:
            p: The created paragraph
        """
        p = text_frame.add_paragraph()
        if is_bullet:
            bullet_text = f"• {text}"
        else:
            bullet_text = text
            
        p.text = bullet_text
        p.alignment = PP_ALIGN.LEFT
        p.level = level
        p.space_after = Pt(5)
        p.line_spacing = 1
        
        run = p.runs[0]
        run.font.size = size
        if color:
            run.font.color.rgb = color
        run.font.bold = bold
        
        return p


    def create_title_slide(self):
        """
        Create a title slide for the presentation.
        
        Returns:
            slide: The created slide
        """
        slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._apply_slide_background(slide)
        # Add title
        prompt = f'Generate a Tittle according to this text{self.text} and make it in 1-2 words'
        response_text = self.generate_response(prompt)
        title_shape = slide.shapes.title
        title_shape.text = response_text
        
        # Style the title
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = self.theme['header_bg']
                run.font.bold = True
        
        # Add subtitle with date
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = f"Created on {time.strftime('%B %d, %Y')}"
        
        return slide


    def generate_preview_slide(self):
        """
        Generate a preview slide with bullet points related to the topic.
        
        Returns:
            tuple: (title, bullet_points)
        """
        title = "PREVIEW"
        prompt = f"Generate 3 - 4 (8 - 10 words) important bullet points related to {self.text}. Keep each point concise (1-2 short sentences ) and make sure to include at least 3 points. Do not include any other text or explanation."
        # Generate response using the AI model
        response_text = self.generate_response(prompt)
        
        # Extract bullet points
        bullet_points = self._extract_bullet_points(response_text)
        
        # Ensure we have at least 3 bullet points
        if len(bullet_points) < 3:
            prompt = f"Generate at least 1 key bullet point from the following text: {self.text}. Cover the main topics related to the bullet points. Only list the topic names; do not provide any explanations or additional text."
            response_text = self.generate_response(prompt)
            bullet_points = self._extract_bullet_points(response_text)
            
            # Fallback if LLM still doesn't provide enough points
            if len(bullet_points) < 3:
                bullet_points = [
                    f"• {point}" for point in bullet_points
                ]
        
        # Create slide
        slide_layout = self.prs.slide_layouts[6]  # Blank slide
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_slide_background(slide)
        # Add header
        self._create_header(slide, title)
        
        # Add bullet points
        text_frame = self._create_text_box(slide)
        
        for i, text in enumerate(bullet_points):
            color = self.theme['bullet_colors'][i % len(self.theme['bullet_colors'])]
            self._add_paragraph(text_frame, text, color=color, is_bullet=True)
            
            # Generate and add a sub-bullet with example
            sub_bullet = f"Example: {text} in real-world applications"
            self._add_paragraph(text_frame, sub_bullet, color=color, level=1, is_bullet=True)
        
        return title, bullet_points

    def generate_aim_slide(self):
        """
        Generate an AIM slide with a concise explanation of the topic.
        
        Returns:
            tuple: (title, response_text)
        """
        title = "AIM"
        prompt = f"Explain {self.text} in one clear, concise sentence."
        response_text = self.generate_response(prompt)
        
        # Create slide
        slide_layout = self.prs.slide_layouts[6]  # Blank slide
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_slide_background(slide)
        # Add header
        self._create_header(slide, title)
        
        # Add centered text
        text_frame = self._create_text_box(slide, centered=True)
        self._add_paragraph(text_frame, response_text, color=self.theme['bullet_colors'][0])
        
        return title, response_text


    def generate_topic_slides(self, bullet_points):
        """
        Generate slides for each bullet point with detailed explanations.
        
        Args:
            bullet_points (list): List of bullet points to create slides for
            
        Returns:
            list: List of (title, content) tuples
        """
        topics = []
        
        for point in bullet_points:
            title = point
            prompt = f"Explain '{point}' related to {self.text} in (2 - 3) bullet points. Keep each point short (1-2 sentences) and make it concise."
            response_text = self.generate_response(prompt)
            
            # Extract bullet points
            sub_bullet_points = self._extract_bullet_points(response_text)
            
            # Ensure we have at least 3 sub-bullet points
            if len(sub_bullet_points) < 3:
                sub_bullet_points = [
                    f"Important aspect of {point}",
                    f"Key development related to {point}",
                    f"Significant impact of {point}"
                ]
            
            topics.append((title, response_text))
            
            # Create slide
            slide_layout = self.prs.slide_layouts[6]  # Blank slide
            slide = self.prs.slides.add_slide(slide_layout)
            self._apply_slide_background(slide)
            # Add header
            self._create_header(slide, title)
            
            # Add bullet points
            text_frame = self._create_text_box(slide)
            
            for i, text in enumerate(sub_bullet_points):
                color = self.theme['bullet_colors'][i % len(self.theme['bullet_colors'])]
                self._add_paragraph(text_frame, text, color=color, is_bullet=True)
                
                # Generate and add a sub-bullet with example
                sub_bullet = f"Example: {text} in real-world applications"
                self._add_paragraph(text_frame, sub_bullet, color=color, level=1, is_bullet=True)
        
        return topics
    
    def generate_conclusion_slide(self):
        """
        Generate a conclusion slide summarizing the topic.
        
        Returns:
            tuple: (title, response_text)
        """
        title = "CONCLUSION"
        prompt = f"Write a concise conclusion (3-4 sentences) for a presentation about {self.text}. and make sure the minimum words (80 - 120) and don't include any other text."
        response_text = self.generate_response(prompt)
        
        # Create slide
        slide_layout = self.prs.slide_layouts[6]  # Blank slide
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_slide_background(slide)
        # Add header
        self._create_header(slide, title)
        
        # Add centered text
        text_frame = self._create_text_box(slide, centered=True)
        self._add_paragraph(text_frame, response_text, color=self.theme['bullet_colors'][0])
        
        return title, response_text

    def generate_thank_you_slide(self, custom_text=None):
        """
        Generate a thank you slide.
        
        Args:
            custom_text (str, optional): Custom text for the thank you slide
            
        Returns:
            str: The title of the slide
        """
        title = custom_text or "JAI HIND"
        
        # Create slide
        slide_layout = self.prs.slide_layouts[6]  # Blank slide
        slide = self.prs.slides.add_slide(slide_layout)
        self._apply_slide_background(slide)
        # Add header
        self._create_header(slide, title)
        
        # Add centered thank you message
        text_frame = self._create_text_box(slide, centered=True)
        thank_you_text = "Developed by JIGYASA\n\nThank you for your attention!"
        self._add_paragraph(text_frame, thank_you_text, 
                           color=self.theme['bullet_colors'][0], 
                           size=Pt(18), 
                           bold=True)
        
        return title

    def generate_presentation(self, output_path=None):
        """
        Generate the complete presentation.
        
        Args:
            output_path (str, optional): Path to save the presentation
            
        Returns:
            str: Path where the presentation was saved
        """
        try:
            print(f"Generating presentation on: {self.text}")
            
            # Create slides
            self.create_title_slide()
            print("✓ Created title slide")
            
            title, bullet_points = self.generate_preview_slide()
            print(f"✓ Created {title} slide")
            
            # title, aim_text = self.generate_aim_slide()
            # print(f"✓ Created {title} slide")
            
            # topics = self.generate_topic_slides(bullet_points)
            # print(f"✓ Created {len(topics)} topic slides")
            
            # title, conclusion = self.generate_conclusion_slide()
            # print(f"✓ Created {title} slide")
            
            # title = self.generate_thank_you_slide()
            # print(f"✓ Created {title} slide")
            
            # Save the presentation
            if not output_path:
                # Create a safe filename from the topic
                safe_topic = "".join(c if c.isalnum() else "_" for c in self.text)
                output_path = f"{safe_topic}_presentation.pptx"
            
            self.prs.save(output_path)
            print(f"✓ Presentation saved to: {output_path}")
            
            return output_path
            
        except Exception as e:
            print(f"Error generating presentation: {str(e)}")
            return None

    def save(self, output_path=None):
        """
        Save the presentation to a file.
        
        Args:
            output_path (str, optional): Path to save the presentation
            
        Returns:
            str: Path where the presentation was saved
        """
        if not output_path:
            # Create a safe filename from the topic
            safe_topic = "".join(c if c.isalnum() else "_" for c in self.text)
            output_path = f"{safe_topic}_presentation.pptx"
        
        try:
            self.prs.save(output_path)
            return output_path
        except Exception as e:
            print(f"Error saving presentation: {str(e)}")
            return None
# Example Usage
if __name__ == "__main__":
    pdf_path = r"G:\Tushar\ppt_generation\pdf\One_Pager_PQC.pdf"  # Replace with actual PDF path
    extractor = PDFExtractor(pdf_path)
    text = extractor.text

    generator = PresentationGenerator(text, model="gemmaModel")

    downloads_dir = os.path.join(os.getcwd(), 'downloads')
    if not os.path.exists(downloads_dir):
        os.makedirs(downloads_dir)
    
    # Generate the presentation in the 'downloads' directory
    output_path = generator.generate_presentation(output_path=os.path.join(downloads_dir, 'presentation.pptx'))

    print(f"Presentation created at: {output_path}")